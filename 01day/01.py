from pptx import Presentation

# 创建一个新的pptx文件
prs = Presentation()

for index,layout in enumerate(prs.slide_layouts):
    print(index, layout.name)
"""
0 Title Slide
1 Title and Content
2 Section Header
3 Two Content
4 Comparison
5 Title Only
6 Blank
7 Content with Caption
8 Picture with Caption
9 Title and Vertical Text
10 Vertical Title and Text
"""

# 获取一个母版
layout = prs.slide_layouts[0]
# 添加一个 幻灯片
slide = prs.slides.add_slide(layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "hello word"
subtitle.text = "python-pptx was here!"

prs.save("01.pptx")




