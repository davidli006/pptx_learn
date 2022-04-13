from pptx import Presentation

prs = Presentation()

tc_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(tc_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Adding a Bullet Slide"

tf = body_shape.text_frame
tf.text = "Find the bullet slide layout"

p = tf.add_paragraph()
p.text = "Use TextFrame for first bullet"
p.level = 1

p = tf.add_paragraph()
p.text = "Use TextFrame.add_paragraph for subsequent bullets"
p.level = 2

p = tf.add_paragraph()
p.text = "Use TextFrame.add_paragraph for 11111"
p.level = 1

p = tf.add_paragraph()
p.text = "Use TextFrame.add_paragraph for 3333"
p.level = 3  # 即使没有第二级,直接设置第三级也是生效的第三级

prs.save("02.pptx")





