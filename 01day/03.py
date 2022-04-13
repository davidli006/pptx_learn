from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

bk_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(bk_layout)

left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = "This is text inside a textbox"

p = tf.add_paragraph()
p.text = "This is second paragraph that's bold"
p.font.bold = True

p2 = tf.add_paragraph()
p2.text = "This is third paragraph that's big"
p2.font.size = Pt(40)

prs.save("03.pptx")
