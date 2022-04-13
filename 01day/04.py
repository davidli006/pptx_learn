from pptx import Presentation
from pptx.util import Inches

img_path = "../media/man.jpg"

prs = Presentation()

bk_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(bk_layout)

left = top = height = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

left = Inches(4)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

width = Inches(3)
slide = prs.slides.add_slide(bk_layout)
movie = slide.shapes.add_movie("../media/video.mp4", left, top,width=width, height=height)

prs.save("04.pptx")