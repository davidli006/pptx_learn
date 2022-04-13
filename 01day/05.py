from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()

title_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_layout)
shapes = slide.shapes

shapes.title.text = "Adding an AutoShape"

left = Inches(0.93)
top = Inches(3)
width = Inches(1.75)
height = Inches(1)

shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
shape.text = "Step 1"

left = left + width - Inches(0.4)
width = Inches(2)

for n in range(2, 6):
    shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.text = f"Step {n}"
    left = left + width - Inches(0.4)

# 所有shapes吗?
for layout in prs.slide_layouts:
    slide = prs.slides.add_slide(layout)

    shapes = slide.shapes
    if shapes.title:  # 母版中都只有一个title
        shapes.title.text = "kakaxi"

prs.save("05.pptx")
