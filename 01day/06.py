from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

to_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(to_layout)

shapes = slide.shapes
shapes.title.text = "Adding a Table"

row = col = 2
left = top = Inches(2)
width = Inches(6)
height = Inches(0.8)

table = shapes.add_table(row, col, left, top, width, height).table

# 设置宽度
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(4)

# 填充数据 第一行
table.cell(0, 0).text = "Foo"
table.cell(0, 1).text = "6"
# 填充数据 第二行
table.cell(1, 0).text = "Baz"
table.cell(1, 1).text = "8"

prs.save("06.pptx")