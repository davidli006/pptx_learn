import os

from pptx import Presentation

def get_all_text(full_path):
    print("-"*10, full_path, "-"*10)
    prs = Presentation(full_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(paragraph.text)


path = "./"
for file in os.listdir(path):
    if file.endswith("py"):
        continue
    full_path = os.path.join(path, file)
    get_all_text(full_path)