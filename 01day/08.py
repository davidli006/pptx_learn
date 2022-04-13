from pptx import Presentation

prs = Presentation()

layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)

lays = prs.slide_master.slide_layouts
lay = lays.get_by_name("Blank")

print(lay.used_by_slides, layout.used_by_slides)
https://python-pptx.readthedocs.io/en/latest/api/slides.html#pptx.slide.SlideMaster


prs.save("08.pptx")
