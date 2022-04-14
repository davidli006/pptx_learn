from pptx import Presentation
from pptx.enum.action import PP_ACTION_TYPE
from pptx.util import Inches

prs = Presentation()

slides = []
for i in range(8):
    layout = prs.slide_layouts[i]
    slides.append(prs.slides.add_slide(layout))


slide_1, slide_5 = slides[1], slides[5]
slide_1.shapes.title.text = "Aim Slide"
slide_5.shapes.title.text = "Enter Button"

img_path = "../media/man.jpg"
left = top = height = Inches(3)
shape = slide_5.shapes.add_picture(img_path, left, top, height=height)

# shape.target_slide = slide_1
# 开始用的直接添加失败, 百度无解, 谷歌才是yyds
shape.click_action.target_slide = slide_1

prs.save("01.pptx")


