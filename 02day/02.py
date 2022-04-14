from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Inches

prs = Presentation()

layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)

shapes = slide.shapes
shapes.title.text = "kakaxi"

shape = shapes[0]

print(shape.adjustments)
print(shape.fill)
fill = shape.fill

# fill.solid()  # 这样填充是黑色

fill.patterned()  # 先调用这个方法就可以对后续填充进行设置

fill.fore_color.rgb = RGBColor(125, 255, 125)
fill.back_color.rgb = RGBColor(20, 10, 20)

# 具体类型在 属性方法 汇总
fill.pattern = MSO_PATTERN_TYPE.DOTTED_DIAMOND

print(shape.shape_type)
print("是否含有文本框", shape.has_text_frame)
print("标题内容", shape.text)
if shape.has_text_frame:
    tf = shape.text_frame

shape = shapes[1]
tf = shape.text_frame
paragraph = tf.add_paragraph()
paragraph.text = "11111"
paragraph.level = 1
paragraph = tf.add_paragraph()
paragraph.text = "2222"
paragraph.level = 1

paragraph = tf.add_paragraph()
paragraph.text = "3333"
paragraph.level = 2

print(tf.auto_size)
tf.auto_size = MSO_AUTO_SIZE.NONE
tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
tf.margin_top = Inches(2)  # 上边距
tf.margin_left = Inches(2)  # 左边距

for pra in tf.paragraphs:
    print(pra.text)

prs.save("02.pptx")