from pptx import Presentation
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

prs = Presentation()

layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)

shapes = slide.shapes
shapes.title.text = "kakaxi"

shape = shapes[-1]
tf = shape.text_frame
paragraph = tf.add_paragraph()
paragraph.text = "2222"
paragraph.level = 2
paragraph.add_line_break()
paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

paragraph = tf.add_paragraph()
paragraph.text = "11111"
paragraph.level = 1
paragraph.add_line_break()


prs.save("03.pptx")
"https://python-pptx.readthedocs.io/en/latest/api/shapes.html#connector-objects"
