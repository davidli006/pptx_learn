from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.util import Pt, Inches, Cm

prs = Presentation()

layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)


shapes = slide.shapes
shapes.title.text = "Connector"

begin_x, begin_y, end_x, end_y = Inches(1), Inches(2), Inches(3), Inches(4)
connector = shapes.add_connector(MSO_CONNECTOR_TYPE.CURVE, begin_x, begin_y, end_x, end_y)
print(connector.name)
print(connector.begin_x, connector.begin_y)
print(connector.end_x, connector.end_y)
print(connector.shape_type)

left = top = Inches(4)
height = Inches(1)
picture = shapes.add_picture("../media/man.jpg", left, top, height=height)
picture.crop_bottom = 0.75
image = picture.image
print(image.blob)
print(image.content_type)
print(image.dpi)
print(image.ext)
print(image.filename)
print(image.sha1)
print(image.size)

# 图片被切割了,只是展示被切割, 图片原大小并不发生变化
with open(f"01.{image.ext}", "wb") as f:
    f.write(image.blob)

slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
shapes.title.text = "Picture"




prs.save("01.pptx")

